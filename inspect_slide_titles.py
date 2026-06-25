from pptx import Presentation

ppt = Presentation(r"c:\Users\jkina\OneDrive - Northeastern University\Northeastern University\Academics\Year 1\EAI 6020 AI System Technologies\Netflix_AI_Case_Study_v5_Academic.pptx")
for idx, slide in enumerate(ppt.slides, start=1):
    title = "<no text>"
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            title = shape.text.strip().split("\n")[0]
            break
    print(f"{idx:02d}: {title}")
