from pptx import Presentation

ppt = Presentation(r"c:\Users\jkina\OneDrive - Northeastern University\Northeastern University\Academics\Year 1\EAI 6020 AI System Technologies\Netflix_AI_Case_Study_v6_Academic.pptx")
for idx, slide in enumerate(ppt.slides, start=1):
    title = next((sh.text.strip().split("\n")[0] for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()), "<no text>")
    print(f"{idx:02d}: {title}")
