from __future__ import annotations

import re
import tempfile
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

BASE = Path(__file__).resolve().parent
COURSE_DIR = BASE.parent.parent
SRC_PPTX = COURSE_DIR / "Netflix_AI_Case_Study_v5_Academic.pptx"
DST_PPTX = COURSE_DIR / "Netflix_AI_Case_Study_v6_Academic.pptx"

# Reordered slide sequence from the current v5 deck to a rubric-aligned flow.
# 1-based slide indices in the source deck.
NEW_ORDER = [
    1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12,
    16, 17, 18, 23, 15, 13, 14, 19, 20, 21, 22,
]

FOOTER_PATTERN = re.compile(
    r"(EAI6020 · Leading AI Projects\s+\|\s+Netflix Recommendation System\s+\|\s+Module 6 Final Presentation\s+\|\s+)(\d+)")


def reorder_slides(ppt: Presentation, order_1_based: list[int]) -> None:
    sld_id_lst = ppt.slides._sldIdLst
    current_ids = list(sld_id_lst)
    reordered = [current_ids[i - 1] for i in order_1_based]
    for sid in list(sld_id_lst):
        sld_id_lst.remove(sid)
    for sid in reordered:
        sld_id_lst.append(sid)


def renumber_visible_slide_numbers(ppt: Presentation) -> None:
    for new_idx, slide in enumerate(ppt.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue

            # Update footer text that includes the slide number.
            if "Module 6 Final Presentation" in text:
                new_text = FOOTER_PATTERN.sub(rf"\g<1>{new_idx}", text)
                if new_text != text:
                    shape.text = new_text
                    continue

            # Update standalone slide number boxes only when they live in the
            # footer region so we do not touch bullet markers like 1/2/3 lists.
            if text.isdigit() and (shape.left > Inches(8.5) or shape.top > Inches(4.5)):
                shape.text = str(new_idx)


def relabel_solution_section(ppt: Presentation) -> None:
    """Replace generic section 5 labels with explicit narrative titles."""
    title_map = {
        18: ("SECTION 5", "Our Alternative Approach: Foundational Principles"),
        19: ("SECTION 5 (CONT.)", "Our Alternative Approach: Concrete Changes"),
        20: ("SECTION 5 (CONT.)", "Risk vs. Business Value: Why the Controls Matter"),
    }

    for slide_idx, (section_tag, title) in title_map.items():
        slide = ppt.slides[slide_idx - 1]
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text.startswith("SECTION 5"):
                shape.text = f"{section_tag}  |  {title}"
                break


def main() -> None:
    if not SRC_PPTX.exists():
        raise FileNotFoundError(f"Missing source deck: {SRC_PPTX}")

    with tempfile.TemporaryDirectory(prefix="pptx_reorder_") as tmp_dir:
        tmp_src = Path(tmp_dir) / "source_v5.pptx"
        tmp_dst = Path(tmp_dir) / "output_v6.pptx"
        shutil.copy2(SRC_PPTX, tmp_src)

        ppt = Presentation(str(tmp_src))
        reorder_slides(ppt, NEW_ORDER)
        renumber_visible_slide_numbers(ppt)
        relabel_solution_section(ppt)
        ppt.save(str(tmp_dst))

        if DST_PPTX.exists():
            DST_PPTX.unlink()
        shutil.copy2(tmp_dst, DST_PPTX)

    print(f"Saved: {DST_PPTX}")
    print(f"Slides: {len(Presentation(str(DST_PPTX)).slides)}")


if __name__ == "__main__":
    main()
