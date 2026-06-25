"""
Build Netflix_AI_Case_Study_v4_Academic.pptx from v3 + demo results.

New slides injected (after slide 14, before the Risk quadrant):
  Slide A  – Our Proposed System: How It Works
  Slide B  – Benchmark Results: 3-Model Comparison  (metrics table)
  Slide C  – Performance Insights: Charts from Real Experiments
  Slide D  – Live Demo: See It In Action
"""

from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
COURSE_DIR = BASE.parent.parent          # EAI 6020 folder
SRC_PPTX   = COURSE_DIR / "Netflix_AI_Case_Study_v3_Academic.pptx"
DST_PPTX   = COURSE_DIR / "Netflix_AI_Case_Study_v4_Academic.pptx"
OUT_DIR    = BASE / "outputs"
CHART1     = OUT_DIR / "metric_overview.png"
CHART2     = OUT_DIR / "accuracy_diversity_tradeoff.png"
METRICS_CSV = OUT_DIR / "metrics.csv"

# Brand colours matching existing deck
DARK_NAVY   = RGBColor(0x1A, 0x1A, 0x2E)
NETFLIX_RED = RGBColor(0xE5, 0x00, 0x14)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE  = RGBColor(0xE8, 0xF4, 0xFD)
MID_GREY    = RGBColor(0x64, 0x74, 0x87)


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text: str, left, top, width, height,
                 font_size=Pt(14), bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_header(slide, section_tag: str, title: str, slide_num: str) -> None:
    """Mimics the existing section header bar at the top."""
    # Thin red accent bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NETFLIX_RED
    bar.line.fill.background()

    add_text_box(slide, f"{section_tag}  |  {title}",
                 Inches(0.25), Inches(0.08), Inches(8.5), Inches(0.38),
                 font_size=Pt(9), bold=True, color=WHITE)

    # Footer bar
    footer = slide.shapes.add_shape(
        1, Inches(0), Inches(7.3), Inches(10), Inches(0.20),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x1A)
    footer.line.fill.background()

    add_text_box(
        slide,
        f"EAI6020 · Leading AI Projects  |  Netflix Recommendation System  |  Module 6 Final Presentation  |  {slide_num}",
        Inches(0.2), Inches(7.3), Inches(9.5), Inches(0.18),
        font_size=Pt(7), color=MID_GREY,
    )


def bullet_box(slide, items: list[tuple[str, str]], left, top, width, height,
               label_size=Pt(11), body_size=Pt(9.5)):
    """Adds a box with bold label + normal body pairs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for label, body in items:
        if not first:
            p = tf.add_paragraph()
            p.space_before = Pt(5)
        else:
            p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = label_size
        run.font.color.rgb = NETFLIX_RED

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = body_size
        r2.font.color.rgb = WHITE
        first = False


# ──────────────────────────────────────────────────────────────────────────────
# Slide factories
# ──────────────────────────────────────────────────────────────────────────────

def make_slide_system(prs: Presentation) -> None:
    """Slide: Our Proposed Lightweight Recommender System."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK_NAVY)
    add_header(slide, "SECTION 5 (CONT.)", "Our Proposed System: How It Works", "15")

    add_text_box(slide, "Lightweight Hybrid Engine — Built & Tested",
                 Inches(0.3), Inches(0.52), Inches(9.4), Inches(0.44),
                 font_size=Pt(20), bold=True, color=WHITE)

    add_text_box(slide,
                 "Grounded in: Li et al. (2024), IEEE Computational Intelligence Magazine, 19(2), 78–95.",
                 Inches(0.3), Inches(0.98), Inches(9.4), Inches(0.25),
                 font_size=Pt(8.5), color=MID_GREY)

    bullet_box(slide, [
        ("Dataset", "MovieLens Latest Small · 100,836 ratings · 610 users · 9,724 items · leave-one-out split"),
        ("Model 1 – Popularity Baseline", "Ranks items by interaction count. Simple, fast, zero personalization."),
        ("Model 2 – Item-Based CF", "Cosine similarity on user-item matrix. Personalized but no diversity control."),
        ("Model 3 – Hybrid Lite (proposed)", "Weighted CF + content-based scores → diversity-aware MMR reranking. "
         "Targets both quality and catalog diversity."),
        ("Evaluation", "Precision@10, Recall@10, NDCG@10, Coverage@10, Intra-list Diversity, Novelty, Latency (ms)."),
        ("Validation protocol", "Leave-one-out per user. Held-out item must appear in top-K recommendations."),
    ],
    Inches(0.3), Inches(1.26), Inches(9.4), Inches(5.6))


def make_slide_metrics(prs: Presentation) -> None:
    """Slide: Benchmark Results table."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK_NAVY)
    add_header(slide, "SECTION 5 (CONT.)", "Benchmark Results: 3-Model Comparison", "16")

    add_text_box(slide, "Evaluation Results  (Top-K = 10, n = 610 users)",
                 Inches(0.3), Inches(0.52), Inches(9.4), Inches(0.40),
                 font_size=Pt(18), bold=True, color=WHITE)

    # Read real metrics
    rows = []
    with METRICS_CSV.open() as f:
        reader = csv.DictReader(f)
        for r in reader:
            rows.append(r)

    # Sort: hybrid first
    order = {"hybrid_lite": 0, "item_cf": 1, "popularity": 2}
    rows.sort(key=lambda r: order.get(r["model"], 9))

    cols = ["model", "ndcg_at_k", "precision_at_k", "recall_at_k",
            "coverage_at_k", "intra_list_diversity", "novelty", "latency_ms"]
    labels = ["Model", "NDCG@10", "Precision@10", "Recall@10",
              "Coverage@10", "Diversity", "Novelty", "Latency (ms)"]

    col_widths = [Inches(1.55), Inches(0.90), Inches(0.90), Inches(0.90),
                  Inches(0.90), Inches(0.90), Inches(0.80), Inches(0.95)]
    total_w = sum(col_widths)
    left_start = (Inches(10) - total_w) / 2
    top_header = Inches(1.02)
    row_h = Inches(0.40)

    # Header row
    x = left_start
    for label, cw in zip(labels, col_widths):
        rect = slide.shapes.add_shape(1, x, top_header, cw, Inches(0.38))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NETFLIX_RED
        rect.line.fill.background()
        txb = slide.shapes.add_textbox(x + Inches(0.04), top_header + Inches(0.04),
                                       cw - Inches(0.08), Inches(0.32))
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = Pt(8.5)
        run.font.color.rgb = WHITE
        x += cw

    # Data rows
    row_colors = [RGBColor(0x14, 0x2A, 0x4A), RGBColor(0x0F, 0x1E, 0x38)]
    for ri, row in enumerate(rows):
        y = top_header + Inches(0.38) + ri * row_h
        x = left_start
        for ci, (col, cw) in enumerate(zip(cols, col_widths)):
            rect = slide.shapes.add_shape(1, x, y, cw, row_h - Inches(0.03))
            rect.fill.solid()
            rect.fill.fore_color.rgb = row_colors[ri % 2]
            rect.line.fill.background()

            val = row[col]
            if col not in ("model", "users_evaluated"):
                try:
                    val = f"{float(val):.4f}"
                except ValueError:
                    pass

            is_highlight = (row["model"] == "hybrid_lite" and ci > 0)
            txb = slide.shapes.add_textbox(x + Inches(0.04), y + Inches(0.06),
                                           cw - Inches(0.08), row_h - Inches(0.1))
            tf = txb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(8.5)
            run.font.bold = is_highlight
            run.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00) if is_highlight else WHITE
            x += cw

    # Lift callout box
    hybrid = next(r for r in rows if r["model"] == "hybrid_lite")
    cf     = next(r for r in rows if r["model"] == "item_cf")
    pop    = next(r for r in rows if r["model"] == "popularity")

    ndcg_lift_cf  = (float(hybrid["ndcg_at_k"]) - float(cf["ndcg_at_k"]))  / float(cf["ndcg_at_k"])  * 100
    ndcg_lift_pop = (float(hybrid["ndcg_at_k"]) - float(pop["ndcg_at_k"])) / float(pop["ndcg_at_k"]) * 100
    lat_hybrid    = float(hybrid["latency_ms"])
    lat_cf        = float(cf["latency_ms"])

    callout = (
        f"Hybrid Lite lifts NDCG@10 by +{ndcg_lift_cf:.1f}% vs Item-CF and +{ndcg_lift_pop:.1f}% vs Popularity  |  "
        f"Latency: {lat_hybrid:.1f} ms vs {lat_cf:.1f} ms (Item-CF)  |  "
        "Diversity & Novelty both highest among personalised models."
    )
    callout_box = slide.shapes.add_shape(1, Inches(0.3), Inches(5.60), Inches(9.4), Inches(0.50))
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
    callout_box.line.color.rgb = NETFLIX_RED
    add_text_box(slide, callout,
                 Inches(0.35), Inches(5.64), Inches(9.3), Inches(0.44),
                 font_size=Pt(8.5), color=WHITE, align=PP_ALIGN.CENTER)

    # Source note
    add_text_box(slide,
                 "Source: Authors' implementation. Dataset: MovieLens Latest Small (GroupLens). "
                 "Paper basis: Li et al. (2024), DOI: 10.1109/MCI.2024.3363984.",
                 Inches(0.3), Inches(6.18), Inches(9.4), Inches(0.28),
                 font_size=Pt(7.5), color=MID_GREY)


def make_slide_charts(prs: Presentation) -> None:
    """Slide: Performance Insights with the two saved charts."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK_NAVY)
    add_header(slide, "SECTION 5 (CONT.)", "Performance Insights: Charts from Real Experiments", "17")

    add_text_box(slide, "Visual Analytics from Our Recommender Engine",
                 Inches(0.3), Inches(0.52), Inches(9.4), Inches(0.38),
                 font_size=Pt(18), bold=True, color=WHITE)

    if CHART1.exists():
        slide.shapes.add_picture(str(CHART1), Inches(0.2), Inches(1.0), Inches(9.6), Inches(3.4))

    if CHART2.exists():
        slide.shapes.add_picture(str(CHART2), Inches(2.0), Inches(4.42), Inches(6.0), Inches(2.65))

    add_text_box(slide,
                 "Figure 3. NDCG@10, Intra-list Diversity, and Inference Latency across all three models.",
                 Inches(0.3), Inches(4.42), Inches(9.4), Inches(0.22),
                 font_size=Pt(7.5), color=MID_GREY, align=PP_ALIGN.CENTER)


def make_slide_demo(prs: Presentation) -> None:
    """Slide: Live Demo callout."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK_NAVY)
    add_header(slide, "LIVE DEMO", "Interactive Recommender System", "18")

    add_text_box(slide, "See It In Action",
                 Inches(0.3), Inches(0.52), Inches(9.4), Inches(0.50),
                 font_size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "streamlit run app.py   →   http://localhost:8501",
                 Inches(1.5), Inches(1.1), Inches(7.0), Inches(0.40),
                 font_size=Pt(13), bold=True, color=NETFLIX_RED, align=PP_ALIGN.CENTER)

    bullet_box(slide, [
        ("Select any user ID",
         "Choose from 610 real MovieLens users. The engine loads their history from the training set."),
        ("Switch models",
         "Toggle between Popularity, Item-CF, and Hybrid Lite. Watch how the ranked list changes."),
        ("See real metrics",
         "NDCG@10, Coverage@10, Diversity, Novelty, and per-query latency are displayed live."),
        ("Accuracy–Diversity tradeoff chart",
         "Visual confirms Hybrid Lite sits at the Pareto frontier: best NDCG and strong diversity."),
        ("Responsible AI checks built in",
         "Diversity-aware MMR reranking prevents filter bubbles at query time with no extra cost."),
    ],
    Inches(0.5), Inches(1.65), Inches(9.0), Inches(5.0),
    label_size=Pt(11), body_size=Pt(9.5))


# ──────────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation(SRC_PPTX)

    # The v3 deck has slides 1-18 (0-indexed 0-17).
    # We want to insert after slide index 13 (= slide 14, "Concrete Changes").
    # python-pptx does not support insert; we add at end then splice via XML.

    insert_after_idx = 13   # 0-based; slide 14 in 1-based numbering

    make_slide_system(prs)
    make_slide_metrics(prs)
    make_slide_charts(prs)
    make_slide_demo(prs)

    # --- Splice: move the 4 new slides (appended at end) to insert_after_idx+1 ---
    from pptx.oxml.ns import qn
    from lxml import etree

    xml_slides = prs.slides._sldIdLst
    all_slide_ids = list(xml_slides)
    n_existing = len(all_slide_ids) - 4          # original slide count before we appended

    # Indexes of newly appended slides in the _sldIdLst
    new_ids = all_slide_ids[n_existing:]

    # Remove new slides from end
    for sid in new_ids:
        xml_slides.remove(sid)

    # Re-insert after the target
    reference_node = all_slide_ids[insert_after_idx]
    for i, sid in enumerate(new_ids):
        reference_node.addnext(sid)
        reference_node = sid     # chain them in order

    prs.save(DST_PPTX)
    print(f"Saved: {DST_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
