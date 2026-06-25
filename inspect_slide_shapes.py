from pptx import Presentation

ppt = Presentation(r"c:\Users\jkina\OneDrive - Northeastern University\Northeastern University\Academics\Year 1\EAI 6020 AI System Technologies\Netflix_AI_Case_Study_v5_Academic.pptx")
for slide_idx in [4, 12, 15, 16, 23]:
    slide = ppt.slides[slide_idx - 1]
    print(f"--- slide {slide_idx} ---")
    for i, shape in enumerate(slide.shapes):
        text = getattr(shape, "text", "")
        if text.strip():
            first = text.strip().split("\n")[0]
            print(i, round(shape.left / 914400, 2), round(shape.top / 914400, 2), round(shape.width / 914400, 2), round(shape.height / 914400, 2), repr(first))
