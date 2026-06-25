from __future__ import annotations

import json
import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
COURSE_DIR = BASE.parent.parent
SRC_PPTX = (COURSE_DIR / "Netflix_AI_Case_Study_v4_Academic.pptx").resolve()
DST_PPTX = (COURSE_DIR / "Netflix_AI_Case_Study_v5_Academic.pptx").resolve()
SUMMARY_JSON = BASE / "outputs" / "summary.json"

DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
NETFLIX_RED = RGBColor(0xE5, 0x00, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MID_GREY = RGBColor(0x64, 0x74, 0x87)


def add_text_box(slide, text: str, left, top, width, height, font_size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return tx_box


def add_header(slide, section_tag: str, title: str, slide_num: str) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NETFLIX_RED
    bar.line.fill.background()

    add_text_box(
        slide,
        f"{section_tag}  |  {title}",
        Inches(0.25),
        Inches(0.08),
        Inches(8.7),
        Inches(0.38),
        font_size=Pt(9),
        bold=True,
    )

    footer = slide.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(10), Inches(0.20))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x1A)
    footer.line.fill.background()

    add_text_box(
        slide,
        f"EAI6020 · Leading AI Projects  |  Netflix Recommendation System  |  Module 6 Final Presentation  |  {slide_num}",
        Inches(0.2),
        Inches(7.3),
        Inches(9.6),
        Inches(0.18),
        font_size=Pt(7),
        color=MID_GREY,
    )


def bullet_block(slide, title: str, body: str, left, top, width, height) -> None:
    add_text_box(slide, title, left, top, width, Inches(0.30), font_size=Pt(12), bold=True, color=NETFLIX_RED)
    add_text_box(slide, body, left, top + Inches(0.28), width, height - Inches(0.28), font_size=Pt(9.8), color=WHITE)


def main() -> None:
    if not SRC_PPTX.exists():
        raise FileNotFoundError(f"Missing source deck: {SRC_PPTX}")
    if not SUMMARY_JSON.exists():
        raise FileNotFoundError(f"Missing summary file: {SUMMARY_JSON}")

    summary = json.loads(SUMMARY_JSON.read_text(encoding="utf-8"))
    opt = summary.get("optimization", {})
    params = opt.get("hybrid_params", {})
    meta = opt.get("metadata", {})
    best = summary.get("best_model", {})

    with tempfile.TemporaryDirectory(prefix="pptx_stage_") as tmp_dir:
        tmp_src = Path(tmp_dir) / "source_v4.pptx"
        tmp_dst = Path(tmp_dir) / "output_v5.pptx"
        shutil.copy2(SRC_PPTX, tmp_src)

        prs = Presentation(str(tmp_src))
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_NAVY

        slide_num = str(len(prs.slides))
        add_header(slide, "SECTION 5 (CONT.)", "Cold-Start Strategy + PSO Optimization", slide_num)

        add_text_box(
            slide,
            "How We Resolved Cold-Start and Tuned the System",
            Inches(0.3), Inches(0.52), Inches(9.4), Inches(0.44),
            font_size=Pt(20), bold=True,
        )

        left_x = Inches(0.35)
        right_x = Inches(5.05)

        bullet_block(
            slide,
            "Cold-start design",
            "For users with little/no history, Hybrid Lite dynamically shifts weight from collaborative signals to content-based similarity and a popularity prior. "
            "This stabilizes early-session relevance while still preserving diversity.",
            left_x,
            Inches(1.12),
            Inches(4.45),
            Inches(1.65),
        )

        bullet_block(
            slide,
            "PSO infusion (where applicable)",
            "Particle Swarm Optimization tuned five production parameters: alpha, diversity_weight, top_pool_size, cold_start_threshold, and popularity_prior_max. "
            "Objective blended ranking quality, coverage, diversity, and latency.",
            right_x,
            Inches(1.12),
            Inches(4.55),
            Inches(1.65),
        )

        p_lines = [
            f"alpha = {params.get('alpha', 0):.3f}",
            f"diversity_weight = {params.get('diversity_weight', 0):.3f}",
            f"top_pool_size = {int(params.get('top_pool_size', 0))}",
            f"cold_start_threshold = {int(params.get('cold_start_threshold', 0))}",
            f"popularity_prior_max = {params.get('popularity_prior_max', 0):.3f}",
        ]
        tuned_text = "\n".join(p_lines)

        panel = slide.shapes.add_shape(1, Inches(0.35), Inches(2.96), Inches(4.45), Inches(2.25))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x3A)
        panel.line.color.rgb = NETFLIX_RED
        add_text_box(slide, "PSO-Tuned Parameters", Inches(0.55), Inches(3.05), Inches(4.0), Inches(0.28), font_size=Pt(12), bold=True, color=WHITE)
        add_text_box(slide, tuned_text, Inches(0.55), Inches(3.38), Inches(4.0), Inches(1.7), font_size=Pt(10), color=WHITE)

        outcome = (
            f"Best model: {best.get('model', 'hybrid_lite')}\n"
            f"NDCG@10: {best.get('ndcg_at_k', 0):.4f}  |  Recall@10: {best.get('recall_at_k', 0):.4f}\n"
            f"Coverage@10: {best.get('coverage_at_k', 0):.4f}  |  Diversity: {best.get('intra_list_diversity', 0):.4f}\n"
            f"Latency: {best.get('latency_ms', 0):.2f} ms"
        )
        panel2 = slide.shapes.add_shape(1, Inches(5.05), Inches(2.96), Inches(4.55), Inches(2.25))
        panel2.fill.solid()
        panel2.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
        panel2.line.color.rgb = NETFLIX_RED
        add_text_box(slide, "Measured Outcome", Inches(5.25), Inches(3.05), Inches(4.15), Inches(0.28), font_size=Pt(12), bold=True, color=WHITE)
        add_text_box(slide, outcome, Inches(5.25), Inches(3.38), Inches(4.15), Inches(1.7), font_size=Pt(10), color=WHITE)

        add_text_box(
            slide,
            f"PSO setup: particles={meta.get('n_particles', 'n/a')}, iterations={meta.get('n_iters', 'n/a')}, eval_users={meta.get('eval_users', 'n/a')}.",
            Inches(0.35), Inches(5.45), Inches(9.3), Inches(0.28), font_size=Pt(8.5), color=MID_GREY, align=PP_ALIGN.CENTER,
        )

        add_text_box(
            slide,
            "Takeaway: we now have an explicit cold-start path and an optimized parameter policy, not just a static recommender design.",
            Inches(0.35), Inches(5.78), Inches(9.3), Inches(0.48), font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )

        prs.save(str(tmp_dst))
        shutil.copy2(tmp_dst, DST_PPTX)

    print(f"Saved: {DST_PPTX}")


if __name__ == "__main__":
    main()
